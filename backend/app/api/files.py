"""File management API routes for agent workspaces."""

import base64
import csv
import io
import mimetypes
import os
import uuid
from pathlib import Path

import aiofiles
from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import FileResponse, Response
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from pydantic import BaseModel

from app.config import get_settings
from app.core.permissions import check_agent_access
from app.core.security import get_current_user
from app.database import get_db
from app.models.user import User
from app.models.workspace import WorkspaceFileRevision
from app.services.workspace_collaboration import (
    acquire_edit_lock,
    content_hash,
    list_revisions,
    read_text_if_exists,
    record_revision,
    release_edit_lock,
    write_workspace_file,
)
from app.services.workspace_paths import WorkspacePathError, resolve_agent_visible_path
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

settings = get_settings()
router = APIRouter(prefix="/agents/{agent_id}/files", tags=["files"])


class FileInfo(BaseModel):
    name: str
    path: str
    is_dir: bool
    size: int = 0
    modified_at: str = ""
    url: str | None = None


class FileContent(BaseModel):
    path: str
    content: str


class FileWrite(BaseModel):
    content: str
    autosave: bool = False
    session_id: str | None = None


class FileLockBody(BaseModel):
    path: str
    session_id: str | None = None


class RestoreRevisionBody(BaseModel):
    revision_id: uuid.UUID


def _agent_base_dir(agent_id: uuid.UUID) -> Path:
    return Path(settings.AGENT_DATA_DIR) / str(agent_id)


def _safe_path(agent_id: uuid.UUID, rel_path: str) -> Path:
    """Ensure the path is within the agent's directory (no path traversal)."""
    base = _agent_base_dir(agent_id)
    full = (base / rel_path).resolve()
    if not str(full).startswith(str(base.resolve())):
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Path traversal not allowed")
    return full


def _visible_path(agent_id: uuid.UUID, rel_path: str, tenant_id: uuid.UUID | None) -> tuple[Path, Path, bool]:
    """Resolve an agent-visible path, including virtual enterprise_info/."""
    try:
        resolved = resolve_agent_visible_path(
            _agent_base_dir(agent_id),
            rel_path,
            workspace_root=Path(settings.AGENT_DATA_DIR),
            tenant_id=str(tenant_id) if tenant_id else None,
        )
    except WorkspacePathError as exc:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=str(exc)) from exc
    return resolved.path, resolved.relative_root, resolved.is_enterprise


@router.get("/", response_model=list[FileInfo])
async def list_files(
    agent_id: uuid.UUID,
    path: str = "",
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """List files and directories in an agent's file system."""
    await check_agent_access(db, current_user, agent_id)
    target, base_abs, is_enterprise = _visible_path(agent_id, path, current_user.tenant_id)
    if is_enterprise:
        target.mkdir(parents=True, exist_ok=True)

    if not target.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Path not found")
    if not target.is_dir():
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Path is not a directory")

    items = []
    base_abs = base_abs.resolve()
    if not path and current_user.tenant_id:
        enterprise_root = (Path(settings.AGENT_DATA_DIR) / f"enterprise_info_{current_user.tenant_id}").resolve()
        enterprise_root.mkdir(parents=True, exist_ok=True)
        items.append(FileInfo(
            name="enterprise_info",
            path="enterprise_info",
            is_dir=True,
            size=0,
            modified_at=str(enterprise_root.stat().st_mtime),
            url=None,
        ))
    for entry in sorted(target.iterdir(), key=lambda e: (not e.is_dir(), e.name)):
        if entry.name == '.gitkeep':
            continue
        if not path and entry.name == "enterprise_info":
            continue
        rel = str(entry.resolve().relative_to(base_abs))
        if is_enterprise:
            rel = f"enterprise_info/{rel}" if rel != "." else "enterprise_info"
        stat = entry.stat()
        items.append(FileInfo(
            name=entry.name,
            path=rel,
            is_dir=entry.is_dir(),
            size=stat.st_size if entry.is_file() else 0,
            modified_at=str(stat.st_mtime),
            url=f"/api/agents/{agent_id}/files/download?path={rel}" if not entry.is_dir() else None
        ))
    return items


@router.get("/content", response_model=FileContent)
async def read_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Read the content of a file."""
    await check_agent_access(db, current_user, agent_id)
    target, _, _ = _visible_path(agent_id, path, current_user.tenant_id)

    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    try:
        async with aiofiles.open(target, "r", encoding="utf-8") as f:
            content = await f.read()
        return FileContent(path=path, content=content)
    except UnicodeDecodeError:
        return FileContent(path=path, content=f"[二进制文件: {target.name}, {target.stat().st_size} bytes]")


def _file_kind(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext == ".csv":
        return "csv"
    if ext in {".html", ".htm"}:
        return "html"
    if ext == ".pdf":
        return "pdf"
    if ext in {".xlsx", ".xls"}:
        return "xlsx"
    if ext in {".docx", ".doc"}:
        return "docx"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    if ext in {".txt", ".log", ".json"}:
        return "text"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    return "binary"


def _find_companion_text_preview(target: Path) -> Path | None:
    for suffix in (".md", ".txt"):
        candidate = target.with_suffix(suffix)
        if candidate.exists() and candidate.is_file():
            return candidate
    return None


def _extract_document_text(target: Path, kind: str) -> str:
    """Best-effort rich document text extraction for lightweight previews."""
    try:
        if kind == "xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(target, read_only=True, data_only=True)
            sheets: list[str] = []
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(max_row=80, max_col=20, values_only=True):
                    rows.append("\t".join("" if cell is None else str(cell) for cell in row))
                sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        if kind == "docx":
            from docx import Document

            doc = Document(str(target))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if kind == "pptx":
            from pptx import Presentation

            prs = Presentation(str(target))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                slides.append(f"Slide {idx}\n" + "\n".join(texts))
            return "\n\n".join(slides)
    except ImportError as exc:
        return f"Missing preview dependency: {exc}"
    except Exception as exc:
        return f"Preview extraction failed: {str(exc)[:200]}"
    return ""


def _detect_csv_delimiter(text: str) -> str:
    lines = [line.strip() for line in text.splitlines() if line.strip()][:10]
    if not lines:
        return ","
    candidates = [",", "，", ";", "\t", "|"]
    scores = {
        candidate: sum(line.count(candidate) for line in lines)
        for candidate in candidates
    }
    return max(scores, key=scores.get) if any(scores.values()) else ","


def _parse_csv_rows(text: str) -> list[list[str]]:
    delimiter = _detect_csv_delimiter(text)
    rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
    normalized: list[list[str]] = []
    for row in rows[:500]:
        values = list(row)
        while values and not str(values[-1] or "").strip():
            values.pop()
        if values:
            normalized.append(values)
    return normalized


@router.get("/preview")
async def preview_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Return a browser-friendly preview payload for Workspace files."""
    await check_agent_access(db, current_user, agent_id)
    target, _, _ = _visible_path(agent_id, path, current_user.tenant_id)
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    kind = _file_kind(path)
    mime_type = mimetypes.guess_type(target.name)[0] or "application/octet-stream"
    download_url = f"/api/agents/{agent_id}/files/download?path={path}"

    if kind in {"markdown", "html", "text"}:
        content = await read_text_if_exists(target)
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "content": content or "",
            "content_hash": content_hash(content or ""),
            "download_url": download_url,
        }
    if kind == "csv":
        content = await read_text_if_exists(target) or ""
        rows = _parse_csv_rows(content)
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "content": content,
            "content_hash": content_hash(content),
            "rows": rows[:500],
            "download_url": download_url,
        }
    if kind == "pdf":
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "url": download_url,
            "download_url": download_url,
        }
    if kind == "xlsx":
        try:
            from openpyxl import load_workbook

            wb = load_workbook(target, read_only=True, data_only=True)
            sheets = []
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(max_row=120, max_col=30, values_only=True):
                    values = ["" if cell is None else str(cell) for cell in row]
                    while values and not str(values[-1] or "").strip():
                        values.pop()
                    if any(value.strip() for value in values):
                        rows.append(values)
                sheets.append({
                    "title": ws.title,
                    "rows": rows,
                })
            wb.close()
            return {
                "path": path,
                "kind": kind,
                "mime_type": mime_type,
                "text": _extract_document_text(target, kind),
                "sheets": sheets,
                "download_url": download_url,
            }
        except Exception as exc:
            return {
                "path": path,
                "kind": kind,
                "mime_type": mime_type,
                "text": f"Preview extraction failed: {str(exc)[:200]}",
                "download_url": download_url,
            }
    if kind in {"docx", "pptx"}:
        extracted_text = _extract_document_text(target, kind)
        companion = _find_companion_text_preview(target)
        companion_content = await read_text_if_exists(companion) if companion is not None else None
        return {
            "path": path,
            "kind": kind,
            "mime_type": mime_type,
            "text": companion_content or extracted_text,
            "companion_path": str(companion.resolve().relative_to(_agent_base_dir(agent_id).resolve())) if companion is not None and not path.startswith("enterprise_info") else None,
            "download_url": download_url,
        }

    companion = _find_companion_text_preview(target)
    if companion is not None:
        content = await read_text_if_exists(companion)
        return {
            "path": path,
            "kind": "text",
            "mime_type": "text/markdown" if companion.suffix.lower() == ".md" else "text/plain",
            "content": content or "",
            "content_hash": content_hash(content or ""),
            "companion_path": str(companion.resolve().relative_to(_agent_base_dir(agent_id).resolve())) if not path.startswith("enterprise_info") else None,
            "download_url": download_url,
        }

    raw = target.read_bytes()
    encoded = base64.b64encode(raw[:1024 * 1024]).decode("ascii")
    return {
        "path": path,
        "kind": kind,
        "mime_type": mime_type,
        "size": target.stat().st_size,
        "base64_sample": encoded,
        "download_url": download_url,
    }


@router.get("/download")
async def download_file(
    agent_id: uuid.UUID,
    path: str,
    token: str = "",
    inline: bool = False,
    credentials: HTTPAuthorizationCredentials | None = Depends(HTTPBearer(auto_error=False)),
    db: AsyncSession = Depends(get_db),
):
    """Download / serve a file from the agent workspace (browser-friendly).
    
    Auth via Bearer header OR `token` query parameter (for <img> tags).
    """
    from app.core.security import decode_access_token

    # Resolve JWT token from either Bearer header or query param
    jwt_token = None
    if credentials:
        jwt_token = credentials.credentials
    elif token:
        jwt_token = token

    if not jwt_token:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Authentication required")

    payload = decode_access_token(jwt_token)
    user_id = payload.get("sub")
    if not user_id:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")

    result = await db.execute(select(User).where(User.id == uuid.UUID(user_id)))
    user = result.scalar_one_or_none()
    if not user or not user.is_active:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="User not found or inactive")

    await check_agent_access(db, user, agent_id)
    target, _, _ = _visible_path(agent_id, path, user.tenant_id)
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")
    return FileResponse(
        path=str(target),
        filename=target.name,
        content_disposition_type="inline" if inline else "attachment",
    )


@router.put("/content")
async def write_file(
    agent_id: uuid.UUID,
    path: str,
    data: FileWrite,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Write content to a file (create or overwrite)."""
    await check_agent_access(db, current_user, agent_id)
    if path.startswith("enterprise_info"):
        if current_user.role not in ("platform_admin", "org_admin"):
            raise HTTPException(status_code=403, detail="Only admins can edit enterprise knowledge base")
        if path.strip("/") == "enterprise_info":
            raise HTTPException(status_code=400, detail="Cannot overwrite enterprise_info root")
        target, _, _ = _visible_path(agent_id, path, current_user.tenant_id)
        target.parent.mkdir(parents=True, exist_ok=True)
        async with aiofiles.open(target, "w", encoding="utf-8") as f:
            await f.write(data.content)
        return {"status": "ok", "path": path, "revision_id": None}

    result = await write_workspace_file(
        db,
        agent_id=agent_id,
        base_dir=_agent_base_dir(agent_id),
        path=path,
        content=data.content,
        actor_type="user",
        actor_id=current_user.id,
        operation="autosave" if data.autosave else "write",
        session_id=data.session_id,
        enforce_human_lock=False,
        merge_user_autosave=data.autosave,
    )
    if not result.ok:
        raise HTTPException(status_code=status.HTTP_409_CONFLICT, detail=result.message)
    await db.commit()
    return {"status": "ok", "path": result.path, "revision_id": result.revision_id}


@router.post("/locks")
async def lock_file(
    agent_id: uuid.UUID,
    data: FileLockBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Acquire or refresh a short-lived human editing lock for a file."""
    await check_agent_access(db, current_user, agent_id)
    lock = await acquire_edit_lock(
        db,
        agent_id=agent_id,
        path=data.path,
        user_id=current_user.id,
        session_id=data.session_id,
    )
    await db.commit()
    return {"status": "ok", "path": lock.path, "expires_at": lock.expires_at.isoformat()}


@router.delete("/locks")
async def unlock_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Release the current user's edit lock for a file."""
    await check_agent_access(db, current_user, agent_id)
    await release_edit_lock(db, agent_id=agent_id, path=path, user_id=current_user.id)
    await db.commit()
    return {"status": "ok", "path": path}


@router.get("/revisions")
async def get_file_revisions(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """List version history for the currently opened Workspace file."""
    await check_agent_access(db, current_user, agent_id)
    if path.startswith("enterprise_info"):
        return []
    revisions = await list_revisions(db, agent_id=agent_id, path=path)
    return [
        {
            "id": str(rev.id),
            "path": rev.path,
            "operation": rev.operation,
            "actor_type": rev.actor_type,
            "actor_id": str(rev.actor_id) if rev.actor_id else None,
            "session_id": rev.session_id,
            "before_content": rev.before_content,
            "after_content": rev.after_content,
            "created_at": rev.created_at.isoformat() if rev.created_at else None,
            "updated_at": rev.updated_at.isoformat() if rev.updated_at else None,
        }
        for rev in revisions
    ]


@router.post("/restore")
async def restore_file_revision(
    agent_id: uuid.UUID,
    data: RestoreRevisionBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Restore a file to a previous revision's after-content."""
    await check_agent_access(db, current_user, agent_id)
    result = await db.execute(
        select(WorkspaceFileRevision).where(
            WorkspaceFileRevision.id == data.revision_id,
            WorkspaceFileRevision.agent_id == agent_id,
        )
    )
    revision = result.scalar_one_or_none()
    if not revision:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Revision not found")
    if revision.after_content is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Cannot restore an empty/deleted revision")

    target = _safe_path(agent_id, revision.path)
    before = await read_text_if_exists(target)
    target.parent.mkdir(parents=True, exist_ok=True)
    async with aiofiles.open(target, "w", encoding="utf-8") as f:
        await f.write(revision.after_content)
    restored = await record_revision(
        db,
        agent_id=agent_id,
        path=revision.path,
        operation="restore",
        actor_type="user",
        actor_id=current_user.id,
        before_content=before,
        after_content=revision.after_content,
    )
    await db.commit()
    return {"status": "ok", "path": revision.path, "revision_id": str(restored.id) if restored else None}


@router.delete("/content")
async def delete_file(
    agent_id: uuid.UUID,
    path: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Delete a file."""
    await check_agent_access(db, current_user, agent_id)
    if path.startswith("enterprise_info") and current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can delete enterprise knowledge base files")
    if path.strip("/") == "enterprise_info":
        raise HTTPException(status_code=400, detail="Cannot delete enterprise_info root")
    target, _, _ = _visible_path(agent_id, path, current_user.tenant_id)

    if not target.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    if target.is_dir():
        import shutil
        shutil.rmtree(target)
    else:
        target.unlink()

    return {"status": "ok", "path": path}


class ImportSkillBody(BaseModel):
    skill_id: str


@router.post("/import-skill")
async def import_skill_to_agent(
    agent_id: uuid.UUID,
    body: ImportSkillBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a global skill into this agent's skills/ workspace folder.

    Copies all files from the global skill registry into
    <agent_workspace>/skills/<folder_name>/.
    """
    await check_agent_access(db, current_user, agent_id)

    from sqlalchemy.orm import selectinload
    from app.models.skill import Skill, SkillFile

    # Load the global skill with its files
    result = await db.execute(
        select(Skill).where(Skill.id == body.skill_id).options(selectinload(Skill.files))
    )
    skill = result.scalar_one_or_none()
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")

    if not skill.files:
        raise HTTPException(status_code=400, detail="Skill has no files")

    # Write each file into the agent's workspace
    base = _agent_base_dir(agent_id)
    skill_dir = base / "skills" / skill.folder_name
    skill_dir.mkdir(parents=True, exist_ok=True)

    written = []
    for f in skill.files:
        file_path = (skill_dir / f.path).resolve()
        # Safety check
        if not str(file_path).startswith(str(base.resolve())):
            continue
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(f.content, encoding="utf-8")
        written.append(f.path)

    return {
        "status": "ok",
        "skill_name": skill.name,
        "folder_name": skill.folder_name,
        "files_written": len(written),
        "files": written,
    }


# Separate router for file uploads (binary) since we need UploadFile
from fastapi import File as FastFile, UploadFile as UploadFileType


upload_router = APIRouter(prefix="/agents/{agent_id}/files", tags=["files"])
DEFAULT_UPLOAD_DIR = "workspace/uploads"


@upload_router.post("/upload")
async def upload_file_to_workspace(
    agent_id: uuid.UUID,
    file: UploadFileType = FastFile(...),
    path: str = "workspace/knowledge_base",
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Upload a binary file to agent workspace."""
    await check_agent_access(db, current_user, agent_id)

    normalized_path = (path or "").strip().strip("/")
    if not normalized_path or normalized_path == ".":
        normalized_path = DEFAULT_UPLOAD_DIR

    # Validate path prefix
    if normalized_path not in {"workspace", "skills"} and not normalized_path.startswith(("workspace/", "skills/")):
        raise HTTPException(status_code=400, detail="右侧根目录视图是 agent 根目录；上传文件时请放到 workspace/ 或 skills/ 目录下")

    base = _agent_base_dir(agent_id)
    target_dir = (base / normalized_path).resolve()
    if not str(target_dir).startswith(str(base.resolve())):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")

    target_dir.mkdir(parents=True, exist_ok=True)
    filename = file.filename or "unnamed"
    # Sanitize filename
    filename = filename.replace("/", "_").replace("\\", "_")
    save_path = target_dir / filename

    content = await file.read()
    save_path.write_bytes(content)

    # Auto-extract text from non-text files
    extracted_path = None
    from app.services.text_extractor import needs_extraction, save_extracted_text
    if needs_extraction(filename):
        txt_file = save_extracted_text(save_path, content, filename)
        if txt_file:
            base_abs = base.resolve()
            extracted_path = str(txt_file.resolve().relative_to(base_abs))

    return {
        "status": "ok",
        "path": f"{normalized_path}/{filename}",
        "url": f"/api/agents/{agent_id}/files/download?path={normalized_path}/{filename}",
        "filename": filename,
        "size": len(content),
        "extracted_text_path": extracted_path,
    }


# ─── Enterprise Knowledge Base ─────────────────────────────────

enterprise_kb_router = APIRouter(prefix="/enterprise/knowledge-base", tags=["enterprise"])


def _enterprise_kb_dir(tenant_id: str) -> Path:
    return Path(settings.AGENT_DATA_DIR) / f"enterprise_info_{tenant_id}" / "knowledge_base"


def _enterprise_info_dir(tenant_id: str) -> Path:
    return Path(settings.AGENT_DATA_DIR) / f"enterprise_info_{tenant_id}"


@enterprise_kb_router.get("/files")
async def list_enterprise_kb_files(
    path: str = "",
    current_user: User = Depends(get_current_user),
):
    """List files in enterprise knowledge base (tenant-scoped)."""
    if not current_user.tenant_id:
        return []
    info_dir = _enterprise_info_dir(str(current_user.tenant_id)).resolve()
    info_dir.mkdir(parents=True, exist_ok=True)

    if path:
        target = (info_dir / path).resolve()
    else:
        target = info_dir
    if not str(target).startswith(str(info_dir)):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")

    if not target.exists() or not target.is_dir():
        return []

    items = []
    for entry in sorted(target.iterdir(), key=lambda e: (not e.is_dir(), e.name)):
        if entry.name == '.gitkeep':
            continue
        rel = str(entry.resolve().relative_to(info_dir.resolve()))
        stat = entry.stat()
        items.append({
            "name": entry.name,
            "path": rel,
            "is_dir": entry.is_dir(),
            "size": stat.st_size if entry.is_file() else 0,
            "url": f"/api/enterprise/knowledge-base/download?path={rel}" if not entry.is_dir() else None
        })
    return items


@enterprise_kb_router.post("/upload")
async def upload_enterprise_kb_file(
    file: UploadFileType = FastFile(...),
    sub_path: str = "",
    current_user: User = Depends(get_current_user),
):
    """Upload a file to enterprise knowledge base (tenant-scoped)."""
    from app.core.security import require_role
    # Only admin can upload to enterprise KB
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can upload to enterprise knowledge base")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    info_dir = _enterprise_info_dir(str(current_user.tenant_id))
    target_dir = (info_dir / sub_path).resolve()
    if not str(target_dir).startswith(str(info_dir.resolve())):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")

    target_dir.mkdir(parents=True, exist_ok=True)
    filename = file.filename or "unnamed"
    filename = filename.replace("/", "_").replace("\\", "_")
    save_path = target_dir / filename

    content = await file.read()
    save_path.write_bytes(content)

    # Auto-extract text from non-text files
    extracted_path = None
    from app.services.text_extractor import needs_extraction, save_extracted_text
    if needs_extraction(filename):
        txt_file = save_extracted_text(save_path, content, filename)
        if txt_file:
            extracted_path = str(txt_file.resolve().relative_to(info_dir.resolve()))

    rel_path = f"{sub_path}/{filename}" if sub_path else filename
    return {
        "status": "ok",
        "path": rel_path,
        "url": f"/api/enterprise/knowledge-base/download?path={rel_path}",
        "filename": filename,
        "size": len(content),
        "extracted_text_path": extracted_path,
    }


@enterprise_kb_router.get("/content")
async def read_enterprise_file(
    path: str,
    current_user: User = Depends(get_current_user),
):
    """Read content of an enterprise knowledge base file (tenant-scoped)."""
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")
    info_dir = _enterprise_info_dir(str(current_user.tenant_id))
    target = (info_dir / path).resolve()
    if not str(target).startswith(str(info_dir.resolve())):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    try:
        content = target.read_text(encoding="utf-8", errors="replace")
        return {"path": path, "content": content}
    except Exception:
        return {"path": path, "content": f"[二进制文件: {target.name}, {target.stat().st_size} bytes]"}


@enterprise_kb_router.put("/content")
async def write_enterprise_file(
    path: str,
    data: FileWrite,
    current_user: User = Depends(get_current_user),
):
    """Write content to an enterprise file (tenant-scoped)."""
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can edit enterprise knowledge base")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    info_dir = _enterprise_info_dir(str(current_user.tenant_id))
    target = (info_dir / path).resolve()
    if not str(target).startswith(str(info_dir.resolve())):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")

    target.parent.mkdir(parents=True, exist_ok=True)
    async with aiofiles.open(target, "w", encoding="utf-8") as f:
        await f.write(data.content)
    return {"status": "ok", "path": path}


@enterprise_kb_router.delete("/content")
async def delete_enterprise_file(
    path: str,
    current_user: User = Depends(get_current_user),
):
    """Delete an enterprise knowledge base file (tenant-scoped)."""
    if current_user.role not in ("platform_admin", "org_admin"):
        raise HTTPException(status_code=403, detail="Only admins can delete enterprise knowledge base files")
    if not current_user.tenant_id:
        raise HTTPException(status_code=400, detail="No tenant associated")

    info_dir = _enterprise_info_dir(str(current_user.tenant_id))
    target = (info_dir / path).resolve()
    if not str(target).startswith(str(info_dir.resolve())):
        raise HTTPException(status_code=403, detail="Path traversal not allowed")
    if not target.exists():
        raise HTTPException(status_code=404, detail="File not found")

    if target.is_dir():
        import shutil
        shutil.rmtree(target)
    else:
        target.unlink()
    return {"status": "ok", "path": path}


# ─── Agent-level ClawHub / URL Skill Import ─────────────────

class ClawhubImportBody(BaseModel):
    slug: str

class UrlImportBody(BaseModel):
    url: str


@router.post("/import-from-clawhub")
async def agent_import_from_clawhub(
    agent_id: uuid.UUID,
    body: ClawhubImportBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a skill from ClawHub directly into this agent's skills/ workspace."""
    await check_agent_access(db, current_user, agent_id)

    from app.api.skills import (
        CLAWHUB_BASE, _fetch_github_directory, _parse_skill_md_frontmatter, _get_github_token,
    )
    import httpx

    slug = body.slug

    # 1. Fetch metadata from ClawHub
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(f"{CLAWHUB_BASE}/v1/skills/{slug}")
            if resp.status_code == 429:
                raise HTTPException(429, "ClawHub rate limit exceeded. Please wait and try again.")
            if resp.status_code != 200:
                raise HTTPException(502, f"ClawHub API error: {resp.status_code}")
            meta = resp.json()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"Failed to connect to ClawHub: {e}")

    skill_info = meta.get("skill", {})
    owner_info = meta.get("owner", {})
    handle = owner_info.get("handle", "").lower()
    if not handle:
        raise HTTPException(400, "Could not determine skill owner from ClawHub metadata")

    # 2. Fetch files from GitHub
    github_path = f"skills/{handle}/{slug}"
    tenant_id = str(current_user.tenant_id) if current_user.tenant_id else None
    token = await _get_github_token(tenant_id)
    files = await _fetch_github_directory("openclaw", "skills", github_path, "main", token)

    # 3. Write to agent workspace: skills/<slug>/
    base = _agent_base_dir(agent_id)
    folder_name = slug
    skill_dir = base / "skills" / folder_name
    skill_dir.mkdir(parents=True, exist_ok=True)

    written = []
    for f in files:
        file_path = (skill_dir / f["path"]).resolve()
        if not str(file_path).startswith(str(base.resolve())):
            continue
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(f["content"], encoding="utf-8")
        written.append(f["path"])

    return {
        "status": "ok",
        "skill_name": skill_info.get("displayName", slug),
        "folder_name": folder_name,
        "files_written": len(written),
        "files": written,
    }


@router.post("/import-from-url")
async def agent_import_from_url(
    agent_id: uuid.UUID,
    body: UrlImportBody,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Import a skill from a GitHub URL directly into this agent's skills/ workspace."""
    await check_agent_access(db, current_user, agent_id)

    from app.api.skills import _parse_github_url, _fetch_github_directory, _get_github_token

    parsed = _parse_github_url(body.url)
    if not parsed:
        raise HTTPException(400, "Invalid GitHub URL")

    owner, repo, branch, path = parsed["owner"], parsed["repo"], parsed["branch"], parsed["path"]
    tenant_id = str(current_user.tenant_id) if current_user.tenant_id else None
    token = await _get_github_token(tenant_id)
    files = await _fetch_github_directory(owner, repo, path, branch, token)
    if not files:
        raise HTTPException(404, "No files found")

    # Derive folder name
    folder_name = path.rstrip("/").split("/")[-1] if path else repo

    # Write to agent workspace
    base = _agent_base_dir(agent_id)
    skill_dir = base / "skills" / folder_name
    skill_dir.mkdir(parents=True, exist_ok=True)

    written = []
    for f in files:
        file_path = (skill_dir / f["path"]).resolve()
        if not str(file_path).startswith(str(base.resolve())):
            continue
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(f["content"], encoding="utf-8")
        written.append(f["path"])

    return {
        "status": "ok",
        "folder_name": folder_name,
        "files_written": len(written),
        "files": written,
    }

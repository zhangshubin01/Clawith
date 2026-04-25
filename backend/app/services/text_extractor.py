"""Extract text from common office file formats.

Supports: PDF, DOCX, XLSX, PPTX
Saves extracted text as a companion .md file alongside the original.
"""

import io
from pathlib import Path

from loguru import logger


# File extensions that need text extraction
EXTRACTABLE_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Text extensions that don't need extraction
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
             ".js", ".ts", ".py", ".html", ".css", ".sh", ".log", ".env"}


def _clean_cell(value: object) -> str:
    text = str(value or "").strip()
    return text.replace("\n", "<br>").replace("|", "\\|")


def _markdown_table(rows: list[list[object]]) -> str:
    cleaned = [[_clean_cell(cell) for cell in row] for row in rows]
    cleaned = [row for row in cleaned if any(cell for cell in row)]
    if not cleaned:
        return ""

    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in cleaned]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def needs_extraction(filename: str) -> bool:
    """Check if a file needs text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in EXTRACTABLE_EXTS


def extract_text(file_bytes: bytes, filename: str) -> str | None:
    """Extract text from a binary file.
    
    Returns extracted text string, or None if extraction fails.
    """
    ext = Path(filename).suffix.lower()
    
    try:
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        elif ext == ".docx":
            return _extract_docx(file_bytes)
        elif ext == ".xlsx":
            return _extract_xlsx(file_bytes)
        elif ext == ".pptx":
            return _extract_pptx(file_bytes)
    except Exception as e:
        logger.error(f"[TextExtractor] Failed to extract from {filename}: {e}")
        return None
    
    return None


def save_extracted_text(save_path: Path, file_bytes: bytes, filename: str) -> Path | None:
    """Extract text and save as a companion .md file.

    For example: report.pdf -> report.md
    Returns the path to the markdown file, or None if extraction failed.
    """
    text = extract_text(file_bytes, filename)
    if not text or not text.strip():
        return None

    md_path = save_path.parent / f"{save_path.stem}.md"
    md_path.write_text(text, encoding="utf-8")
    logger.info(f"[TextExtractor] Extracted {len(text)} chars from {filename} -> {md_path.name}")
    return md_path


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber
    
    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_parts = []
            text = page.extract_text()
            if text and text.strip():
                page_parts.append(text.strip())
            
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                if table:
                    table_md = _markdown_table(table)
                    if table_md:
                        page_parts.append(table_md)

            if page_parts:
                pages.append(f"## 第 {i + 1} 页\n\n" + "\n\n".join(page_parts))
    
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    
    doc = Document(io.BytesIO(data))
    parts = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve heading hierarchy
            style_name = para.style.name if para.style and para.style.name else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                parts.append(f"{'#' * level} {text}")
            elif "List Bullet" in style_name:
                parts.append(f"- {text}")
            elif "List Number" in style_name:
                parts.append(f"1. {text}")
            else:
                parts.append(text)
    
    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        table_md = _markdown_table(rows)
        if table_md:
            parts.append("## 表格\n\n" + table_md)
    
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [c if c is not None else "" for c in row]
            if any(str(c).strip() for c in cells):
                rows.append(cells)
        
        table_md = _markdown_table(rows)
        if table_md:
            parts.append(f"## 工作表: {sheet}\n\n" + table_md)
    
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    
    prs = Presentation(io.BytesIO(data))
    parts = []
    
    for i, slide in enumerate(prs.slides):
        texts = []
        tables = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                table_md = _markdown_table(rows)
                if table_md:
                    tables.append(table_md)
        
        slide_parts = []
        if texts:
            slide_parts.append("\n\n".join(texts))
        slide_parts.extend(tables)
        if slide_parts:
            parts.append(f"## 幻灯片 {i + 1}\n\n" + "\n\n".join(slide_parts))
    
    return "\n\n".join(parts)
